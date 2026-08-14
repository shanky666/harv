# HarvestLenz - Core Four (Mango/Pineapple/Grapes/Pomegranate) 30-slide deck
# Uses the AgroZone brand palette: green 2D5016, orange CC5500, gray-green 5A5F54,
# dark bg 22281C/1A1A1A, light bg F4F6F0/FFFFFF.
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------- palette ----------
GREEN      = RGBColor(0x2D, 0x50, 0x16)
ORANGE     = RGBColor(0xCC, 0x55, 0x00)
GRAYGREEN  = RGBColor(0x5A, 0x5F, 0x54)
DARK       = RGBColor(0x22, 0x28, 0x1C)
DARK2      = RGBColor(0x1A, 0x1A, 0x1A)
LIGHT      = RGBColor(0xF4, 0xF6, 0xF0)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
MID        = RGBColor(0x86, 0x8C, 0x7E)

SW, SH = Inches(13.333), Inches(7.5)

FONT_H = "Calibri Light"
FONT_B = "Calibri"

VIDEO_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "demo_videos")
VIDEO_BY_FRUIT = {
    "Mango":       "mango_demo.mp4",
    "Pineapple":   "pineapple_demo.mp4",
    "Grapes":      "grapes_demo.mp4",
    "Pomegranate": "pomegranate_demo.mp4",
}
VIDEO_POSTER = os.path.join(VIDEO_DIR, "poster_frame.png")

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]

PAGE = [0]

def new_slide(bg):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    return s

def rect(s, x, y, w, h, color, line=False):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    r.fill.solid(); r.fill.fore_color.rgb = color
    if line:
        r.line.color.rgb = color
    else:
        r.line.fill.background()
    r.shadow.inherit = False
    return r

def txt(s, x, y, w, h, text, size, bold=False, color=WHITE, font=FONT_B, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.0, line_spacing=None):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(4)
    tf.margin_top = tf.margin_bottom = Pt(2)
    if isinstance(text, str):
        text = [text]
    first = True
    for i, para in enumerate(text):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(spacing)
        if line_spacing:
            p.line_spacing = line_spacing
        if isinstance(para, tuple):
            runs, psize, pbold, pcolor, pfont, palign = para
            p.alignment = palign
            for rtxt, rsize, rbold, rcolor in runs:
                r = p.add_run(); r.text = rtxt
                r.font.size = Pt(rsize); r.font.bold = rbold
                r.font.color.rgb = rcolor; r.font.name = FONT_B
        else:
            r = p.add_run(); r.text = para
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = color; r.font.name = font
    return tb

def footer(s, num, dark=True):
    c = GRAYGREEN if dark else GRAYGREEN
    rect(s, 0, SH - Inches(0.32), SW, Inches(0.32), DARK2 if dark else LIGHT)
    txt(s, Inches(0.45), SH - Inches(0.30), Inches(4), Inches(0.28), "HarvestLenz  \u00b7  Core Four Deep-Dive", 9, color=MID if dark else GRAYGREEN, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(11.9), SH - Inches(0.30), Inches(1.0), Inches(0.28), f"{num:02d}", 10, bold=True, color=ORANGE, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

def logo(s, x, y, size, color=ORANGE):
    txt(s, x, y, Inches(size), Inches(size), "\u2740", size*28, bold=True, color=color)

def header(s, title, sub=None, num=None):
    if num is not None: PAGE[0] = num
    rect(s, 0, 0, SW, Inches(0.09), GREEN)
    logo(s, Inches(0.45), Inches(0.28), 0.42)
    txt(s, Inches(1.05), Inches(0.26), Inches(11.0), Inches(0.7), title, 34, bold=True, color=GREEN, font=FONT_H)
    if sub:
        txt(s, Inches(1.05), Inches(0.84), Inches(11.6), Inches(0.42), sub, 16, color=GRAYGREEN)
    rect(s, Inches(1.05), Inches(1.32), Inches(0.9), Inches(0.06), ORANGE)

def bullets(s, x, y, w, h, items, size=17, color=GRAYGREEN, gap=10, line_spacing=1.1):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(2); tf.margin_top = tf.margin_bottom = Pt(2)
    first = True
    for it in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap); p.line_spacing = line_spacing
        if isinstance(it, tuple):
            head, body = it
            r = p.add_run(); r.text = head + "  "
            r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = GREEN; r.font.name = FONT_B
            r2 = p.add_run(); r2.text = body
            r2.font.size = Pt(size); r2.font.bold = False; r2.font.color.rgb = color; r2.font.name = FONT_B
        else:
            r = p.add_run(); r.text = it
            r.font.size = Pt(size); r.font.bold = False; r.font.color.rgb = color; r.font.name = FONT_B
    return tb

def card(s, x, y, w, h, title, body, tcolor=GREEN, bg=LIGHT, accent=ORANGE, title_size=17, body_size=15):
    rect(s, x, y, w, h, bg)
    rect(s, x, y, Inches(0.07), h, accent)
    txt(s, x+Inches(0.22), y+Inches(0.12), w-Inches(0.4), Inches(0.38), title, title_size, bold=True, color=tcolor)
    if isinstance(body, list):
        bullets(s, x+Inches(0.22), y+Inches(0.56), w-Inches(0.4), h-Inches(0.66), body, size=body_size, gap=6, line_spacing=1.08)
    else:
        txt(s, x+Inches(0.22), y+Inches(0.56), w-Inches(0.4), h-Inches(0.66), body, body_size, color=GRAYGREEN, line_spacing=1.08)

def section_title(num, kicker, title, sub, items=None):
    s = new_slide(DARK)
    txt(s, Inches(0.6), Inches(0.85), Inches(3), Inches(1.2), num, 72, bold=True, color=ORANGE, font=FONT_H)
    rect(s, Inches(0.65), Inches(2.05), Inches(0.12), Inches(0.75), ORANGE)
    txt(s, Inches(0.95), Inches(2.0), Inches(11.4), Inches(1.2), kicker, 18, bold=True, color=MID)
    txt(s, Inches(0.95), Inches(2.45), Inches(11.6), Inches(1.1), title, 46, bold=True, color=WHITE, font=FONT_H)
    if sub:
        txt(s, Inches(0.95), Inches(3.55), Inches(11.2), Inches(0.8), sub, 18, color=RGBColor(0xC9, 0xCD, 0xC0))
    if items:
        bullets(s, Inches(0.95), Inches(4.6), Inches(11.4), Inches(2.3), items, size=19, color=RGBColor(0xDD, 0xE0, 0xD5))
    footer(s, PAGE[0], dark=True)
    return s

def content(title, sub=None, num=None):
    s = new_slide(LIGHT)
    header(s, title, sub, num)
    return s

def video_placeholder(name, num):
    if num is not None: PAGE[0] = num
    s = new_slide(DARK)
    rect(s, 0, 0, SW, Inches(0.09), GREEN)
    logo(s, Inches(0.45), Inches(0.28), 0.42)
    txt(s, Inches(1.05), Inches(0.26), Inches(11.0), Inches(0.7), f"{name} \u2014 Video Demo", 34, bold=True, color=WHITE, font=FONT_H)
    txt(s, Inches(1.05), Inches(0.84), Inches(11.6), Inches(0.42), "A real run of the model in action.", 16, color=RGBColor(0xC9, 0xCD, 0xC0))
    rect(s, Inches(1.05), Inches(1.32), Inches(0.9), Inches(0.06), ORANGE)
    movie_file = os.path.join(VIDEO_DIR, VIDEO_BY_FRUIT.get(name, "")) if name in VIDEO_BY_FRUIT else ""
    if movie_file and os.path.exists(movie_file):
        poster = VIDEO_POSTER if os.path.exists(VIDEO_POSTER) else None
        s.shapes.add_movie(movie_file, Inches(1.6), Inches(1.8), Inches(10.1), Inches(4.6),
                           poster_frame_image=poster, mime_type="video/mp4")
        txt(s, Inches(1.6), Inches(1.8), Inches(10.1), Inches(0.5),
            "\u25B6  Click to play the demo reel", 15, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
        txt(s, Inches(1.6), Inches(6.62), Inches(10.1), Inches(0.5),
            "Real grading output rendered to H.264 \u2014 1920\u00d71080, 24 fps. Segment \u2192 crop \u2192 grade \u2192 shelf life \u2192 market call.", 14, color=RGBColor(0xC9, 0xCD, 0xC0), align=PP_ALIGN.CENTER)
    else:
        r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.6), Inches(1.8), Inches(10.1), Inches(4.6))
        r.fill.solid(); r.fill.fore_color.rgb = RGBColor(0x28, 0x2E, 0x22)
        r.line.color.rgb = MID; r.line.width = Pt(1.5)
        r.shadow.inherit = False
        txt(s, Inches(1.6), Inches(2.7), Inches(10.1), Inches(1.1), "\u25B6", 60, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
        txt(s, Inches(1.6), Inches(3.8), Inches(10.1), Inches(0.7), f"Insert {name} demo video here", 26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, Inches(1.6), Inches(4.55), Inches(10.1), Inches(0.6), "Segment \u2192 crop \u2192 grade \u2192 shelf life \u2192 market call \u2014 captured on a real basket", 16, color=RGBColor(0xC9, 0xCD, 0xC0), align=PP_ALIGN.CENTER)
        txt(s, Inches(1.6), Inches(5.3), Inches(10.1), Inches(0.5), "Replace this placeholder with an MP4 or a screen-recording link.", 14, color=MID, align=PP_ALIGN.CENTER)
    footer(s, num, dark=True)
    return s

# =============================================================
# SLIDE 1 - Title
# =============================================================
s = new_slide(DARK)
rect(s, 0, 0, SW, Inches(0.18), GREEN)
rect(s, 0, SH - Inches(0.18), SW, Inches(0.18), ORANGE)
logo(s, Inches(6.1), Inches(0.95), 0.9)
txt(s, Inches(1.5), Inches(2.15), Inches(10.3), Inches(1.1), "HARVESTLENZ", 72, bold=True, color=WHITE, font=FONT_H, align=PP_ALIGN.CENTER)
txt(s, Inches(1.5), Inches(3.35), Inches(10.3), Inches(0.7), "Fruit Quality Grading System \u2014 Deep-Dive", 28, color=ORANGE, align=PP_ALIGN.CENTER)
txt(s, Inches(2.0), Inches(4.35), Inches(9.3), Inches(0.7), "The Core Four:  Mango  \u00b7  Pineapple  \u00b7  Grapes  \u00b7  Pomegranate", 26, bold=True, color=RGBColor(0xC9, 0xCD, 0xC0), align=PP_ALIGN.CENTER)
txt(s, Inches(2.5), Inches(5.3), Inches(8.3), Inches(0.7), "One photo in.  Grade, shelf life, and market call out.", 20, color=GRAYGREEN, align=PP_ALIGN.CENTER)
txt(s, Inches(2.5), Inches(6.25), Inches(8.3), Inches(0.6), "Part of the AgroZone Suite  \u00b7  One soil reading, full crop lifecycle", 17, color=MID, align=PP_ALIGN.CENTER)
PAGE[0] = 1
footer(s, 1, dark=True)

# =============================================================
# SLIDE 2 - Agenda
# =============================================================
s = content("What This Deck Covers", "A complete technical walkthrough of the four best-supported fruits.", 2)
cols = [
    ("01  Foundation", ["The problem of manual grading", "Solution at a glance", "Single vs Basket mode", "The four tuned originals"]),
    ("02  The Pipeline", ["Segmentation (classical OpenCV)", "Background removal", "MobileNetV2 grading", "Defect fusion + async flow"]),
    ("03  The Core Four", ["Mango \u2014 overview, data, model", "Pineapple \u2014 overview, data, model", "Grapes \u2014 overview, data, model", "Pomegranate \u2014 overview, data, model"]),
    ("04  Results & Future", ["End-to-end walkthrough", "API reference", "Tech stack", "Roadmap & closing"]),
]
x = Inches(0.55)
for t, items in cols:
    card(s, x, Inches(1.75), Inches(2.95), Inches(4.9), t, items)
    x += Inches(3.12)
footer(s, 2, dark=False)

# =============================================================
# SLIDE 3 - The Problem
# =============================================================
s = content("The Problem", "Manual grading doesn't scale \u2014 and one model can't grade every fruit.", 3)
card(s, Inches(0.55), Inches(1.7), Inches(4.0), Inches(4.6), "\u201cManual grading is inconsistent\u201d", [
    ("Slow & subjective", "sorting by eye varies between graders"),
    ("No audit trail", "decisions can't be traced or defended"),
    ("Labor-bound", "does not scale with harvest volumes"),
])
card(s, Inches(4.8), Inches(1.7), Inches(4.0), Inches(4.6), "\u201cOne classifier underperforms\u201d", [
    ("Damage differs by fruit", "mango rot \u2260 pineapple bruise \u2260 grape shrivel"),
    ("Generic models miss", "a single network can't capture 4 defect types"),
    ("Per-fruit models win", "one weight file per fruit type"),
])
card(s, Inches(9.05), Inches(1.7), Inches(3.7), Inches(4.6), "\u201cNo business signal\u201d", [
    ("Grade alone isn't enough", "sellers also need shelf life"),
    ("No market call", "export vs local vs processing"),
    ("No pricing", "grade \u2192 rupee value per fruit"),
])
footer(s, 3, dark=False)

# =============================================================
# SLIDE 4 - Solution at a glance
# =============================================================
s = content("Our Solution", "One photo in.  Grade, shelf life, and market call out.", 4)
card(s, Inches(0.55), Inches(1.75), Inches(3.6), Inches(4.7), "\u2460 Basket photo + fruit type", [
    ("Known fruit", "user picks from the dropdown \u2014 no classifier needed"),
    ("One upload", "single fruit or a whole basket"),
    ("Async", "job starts instantly, results polled"),
], accent=GREEN)
card(s, Inches(4.5), Inches(1.75), Inches(3.6), Inches(4.7), "\u2461 Segment + clean", [
    ("Classical vision", "HSV thresholds + watershed + contours"),
    ("Mask crops", "background removed before any ML"),
    ("No ML detector", "deterministic and fast"),
], accent=GREEN)
card(s, Inches(8.45), Inches(1.75), Inches(3.6), Inches(4.7), "\u2462 Grade + advise", [
    ("MobileNetV2", "per-fruit 3-class: Better / Good / Reject"),
    ("Shelf life", "days remaining from grade + defect score"),
    ("Market + price", "export / supermarket / processing, rupee value"),
], accent=GREEN)
txt(s, Inches(0.55), Inches(6.6), Inches(12.2), Inches(0.4),
    "Every stage is logged and traceable \u2014 the same glass-box philosophy as the rest of the AgroZone Suite.", 16, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
footer(s, 4, dark=False)

# =============================================================
# SLIDE 6 - Single vs Basket mode
# =============================================================
s = content("Capture Modes: Single vs Basket", "The same pipeline, two capture intents.", 5)
card(s, Inches(0.55), Inches(1.7), Inches(5.9), Inches(4.9), "Single mode  (is_single = true)", [
    ("Whole frame = one fruit", "full-image mask, one grade and one price"),
    ("Best for", "one mango, one pineapple, a single fruit on a plain background"),
    ("Pipeline", "segment \u2192 crop \u2192 grade \u2192 shelf/market once"),
], accent=GREEN)
card(s, Inches(6.75), Inches(1.7), Inches(6.0), Inches(4.9), "Basket mode  (is_single = false)", [
    ("Per-fruit detection", "each fruit segmented and graded individually"),
    ("Best for", "a pile of grapes, a bowl of mangoes, a mixed basket"),
    ("Output", "per-fruit table + basket summary + basket price"),
], accent=ORANGE)
txt(s, Inches(0.55), Inches(6.7), Inches(12.2), Inches(0.4),
    "All 21 supported crops accept both modes in the API \u2014 the Core Four are simply the best-tuned for baskets.", 16, bold=True, color=GRAYGREEN, align=PP_ALIGN.CENTER)
footer(s, 5, dark=False)

# =============================================================
# SLIDE 7 - Section 1 divider + Architecture
# =============================================================
section_title("01", "SECTION 1 \u00b7 THE PIPELINE", "Five stages, from photo to market call.",
    "Every basket runs through the same async pipeline.", [
    ("1", "Segmentation \u2014 OpenCV watershed + contour analysis, no ML detector"),
    ("2", "Background removal \u2014 mask-based crop extraction (224\u00d7224)"),
    ("3", "Per-fruit grading \u2014 independent MobileNetV2 models (Better / Good / Reject)"),
    ("4", "Shelf life + market + pricing \u2014 grade feeds every business output"),
    ("5", "Async processing \u2014 /analyze returns a session id immediately"),
])

# =============================================================
# SLIDE 8 - System Architecture
# =============================================================
s = content("System Architecture", "Classical vision first \u2014 deep learning only where it earns its keep.", 7)
card(s, Inches(0.55), Inches(1.65), Inches(3.4), Inches(4.9), "FastAPI Backend", [
    ("Async, session-based", "/analyze, /analysis/{id}, /analysis/{id}/status, /stats"),
    ("Background jobs", "grading runs off the request thread"),
    ("Storage layer", "uploads / original, processed, crops, reports"),
], accent=GREEN)
card(s, Inches(4.25), Inches(1.65), Inches(3.4), Inches(4.9), "OpenCV Segmentation", [
    ("HSV thresholds", "per-fruit tuned color masks"),
    ("Watershed", "separates touching fruits"),
    ("Contour fallback", "plus full-frame fallback for single mode"),
], accent=GREEN)
card(s, Inches(7.95), Inches(1.65), Inches(3.4), Inches(4.9), "Keras Models", [
    ("MobileNetV2", "ImageNet-frozen base, custom head"),
    ("One file per fruit", "mango / pineapple / grapes / pomegranate .keras"),
    ("3-class softmax", "Better \u00b7 Good \u00b7 Reject"),
], accent=GREEN)
card(s, Inches(11.65), Inches(1.65), Inches(1.15), Inches(4.9), "SQLite", [
    ("Async SQLAlchemy", "sessions + per-fruit results"),
    ("Dashboard stats", "/stats aggregates"),
], accent=GREEN)
txt(s, Inches(0.55), Inches(6.75), Inches(12.2), Inches(0.4),
    "Request flow:  Upload \u2192 OpenCV segmentation \u2192 background removal \u2192 MobileNetV2 grading \u2192 shelf-life + market \u2192 dashboard", 16, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
footer(s, 7, dark=False)

# =============================================================
# SLIDE 9 - Segmentation stage
# =============================================================
s = content("Stage 1 \u2014 Segmentation (Classical OpenCV)", "No ML detector \u2014 deterministic, inspectable, fast.", 8)
bullets(s, Inches(0.55), Inches(1.6), Inches(6.4), Inches(4.6), [
    ("Foreground mask", "S/V thresholding, tuned per fruit type"),
    ("Watershed separation", "distance transform + markers split touching fruits"),
    ("Contour analysis", "external contours as the fallback strategy"),
    ("Multi-strategy scoring", "pick the mask+split with the best valid-area score"),
    ("Full-frame fallback", "used when nothing valid is detected (also single mode)"),
], size=18)
card(s, Inches(7.35), Inches(1.6), Inches(5.4), Inches(4.9), "Tuned color masks (per fruit)", [
    ("Grapes", "green (25\u201385\u00b0H) \u222a purple (120\u2013165\u00b0H)"),
    ("Pomegranate", "red (0\u201310\u00b0H) \u222a red (160\u2013180\u00b0H)"),
    ("Pineapple", "yellow (15\u201335\u00b0H) \u222a green (25\u201385\u00b0H)"),
    ("Mango", "generic S/V threshold (S>60, V>80)"),
], accent=ORANGE)
footer(s, 8, dark=False)

# =============================================================
# SLIDE 10 - Background removal
# =============================================================
s = content("Stage 2 \u2014 Background Removal", "Mask-based crop extraction isolates each fruit before grading.", 9)
bullets(s, Inches(0.55), Inches(1.6), Inches(6.6), Inches(4.6), [
    ("Mask-driven", "uses the segmentation mask, not a box, to cut the fruit out"),
    ("Padding", "5% margin around the bbox so edges are not clipped"),
    ("Target size", "224\u00d7224 \u2014 the CNN input shape"),
    ("Edge refinement", "optional; disabled in the main pipeline for speed"),
    ("Why it matters", "the grader sees only the fruit, so it learns fruit defects, not background noise"),
], size=18)
card(s, Inches(7.45), Inches(1.6), Inches(5.3), Inches(4.9), "What the classifier receives", [
    ("Before", "raw crop: background, surface, shadows, other fruits"),
    ("After", "clean fruit on a neutral background"),
    ("Result", "higher signal for Better / Good / Reject separation"),
], accent=ORANGE)
footer(s, 9, dark=False)

# =============================================================
# SLIDE 11 - The CNN
# =============================================================
s = content("Stage 3 \u2014 The Grading CNN (MobileNetV2)", "One shared architecture, four independent weight files.", 10)
rect(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(2.5), WHITE)
txt(s, Inches(0.8), Inches(1.85), Inches(11.7), Inches(0.45), "Architecture", 19, bold=True, color=GREEN)
stages = ["Image 224\u00d7224\u00d73", "MobileNetV2 base (ImageNet, frozen)", "GlobalAvgPool", "Dense 256 + BN", "Dropout 0.4", "Dense 128", "Dropout 0.3", "Dense 3 softmax"]
N = len(stages)
box_w = 1.28
gap_w = 0.24
start_x = (12.6 - (N * box_w + (N - 1) * gap_w))  # center inside container 0.6..12.7
for i, st in enumerate(stages):
    x = Inches(start_x + i * (box_w + gap_w))
    rect(s, x, Inches(2.4), Inches(box_w), Inches(1.0), GREEN if i < 2 else (ORANGE if i == 7 else RGBColor(0x5A, 0x6B, 0x3A)))
    txt(s, x, Inches(2.44), Inches(box_w), Inches(0.95), st, 11, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if i < N - 1:
        txt(s, Inches(start_x + i * (box_w + gap_w) + box_w), Inches(2.4), Inches(gap_w), Inches(1.0), "\u2192", 12, bold=True, color=ORANGE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
bullets(s, Inches(0.6), Inches(4.45), Inches(12.1), Inches(2.3), [
    ("Input", "background-removed 224\u00d7224 RGB crop, MobileNetV2 preprocess"),
    ("Training", "two-phase transfer learning \u2014 head-only (lr 1e-3, ~10 epochs) then unfreeze from layer 100 (lr 1e-4)"),
    ("Augmentation", "rotation, brightness, zoom, horizontal flip, hue shift"),
    ("Class ordering", "alphabetical for TensorFlow: Better / Good / Reject"),
], size=17)
footer(s, 10, dark=False)

# =============================================================
# SLIDE 12 - Defect fusion + async
# =============================================================
s = content("Stage 4 + 5 \u2014 Defect Fusion, Business Signals, Async", "The CNN is blended with classical defect features \u2014 and the UI never waits.", 11)
card(s, Inches(0.55), Inches(1.6), Inches(6.0), Inches(4.9), "Classical defect features", [
    ("Dark-spot detection", "V-channel threshold per fruit (mango 40, grapes 45, pomegranate 50, pineapple 60)"),
    ("Defect score", "spot area + color uniformity \u2192 0\u20131"),
    ("Defect list", "minor spots / moderate blemishes / severe rot or bruising / shape deformation"),
    ("Override rules", "defect > 0.7 \u2192 Reject; low CNN margin + defect > 0.5 \u2192 Reject"),
    ("Confidence gate", "results below 0.40 confidence are dropped"),
], accent=ORANGE)
card(s, Inches(6.85), Inches(1.6), Inches(5.9), Inches(4.9), "Business signals + async", [
    ("Shelf life", "grade + defect score \u2192 days remaining"),
    ("Market call", "Export / Supermarket / Processing Industry"),
    ("Pricing", "base price \u00d7 quality multiplier (Good 1.0, Better 0.75, Reject 0.40)"),
    ("Async", "/analyze returns session_id instantly"),
    ("Polling", "GET /analysis/{id}/status until complete"),
], accent=GREEN)
footer(s, 11, dark=False)

# =============================================================
# Core Four data
# =============================================================
FRUITS = [
    dict(
        name="Mango",
        num="02",
        overview=[
            ("What is graded", "ripeness, anthracnose, stem-end rot, black spot, bruising, surface blemishes"),
            ("Typical defects", "dark lesions, soft spots, shrivel, over-ripening"),
            ("Why it is hard", "color spans green \u2192 yellow \u2192 orange, so a generic mask under-segments"),
            ("Grade meaning", "Better = exportable, Good = retail, Reject = processing/discard"),
        ],
        data=[
            ("Format", "dataset-mango / Better | Good | Reject"),
            ("3 classes", "alphabetical ordering used by TensorFlow"),
            ("Augmentation", "rotation, brightness, zoom, flip, hue"),
            ("Splits", "train / val (val holds a held-out set)"),
        ],
        model=[
            ("Weights", "mango.keras \u2014 27.6 MB MobileNetV2"),
            ("Head", "GAP \u2192 Dense256-BN \u2192 Dropout0.4 \u2192 Dense128 \u2192 Dropout0.3 \u2192 softmax(3)"),
            ("V-threshold", "40 \u2014 the lowest of the Core Four (dark spots stand out on pale skin)"),
            ("Segmentation", "generic S/V threshold (no dedicated color mask)"),
        ],
    ),
    dict(
        name="Pineapple",
        num="03",
        overview=[
            ("What is graded", "skin color, bruising, mold, sunscald, crown condition"),
            ("Typical defects", "brown or black patches, mold at the base, dried crown"),
            ("Why it is hard", "cylindrical shape + tough rind \u2014 shading and highlights confuse simple masks"),
            ("Grade meaning", "Better = fresh export, Good = local retail, Reject = juicing/discard"),
        ],
        data=[
            ("Format", "dataset-pineapple / Better | Good | Reject"),
            ("3 classes", "alphabetical ordering used by TensorFlow"),
            ("Augmentation", "rotation, brightness, zoom, flip, hue"),
            ("Splits", "train / val with held-out validation"),
        ],
        model=[
            ("Weights", "pineapple.keras \u2014 27.6 MB MobileNetV2"),
            ("Head", "same shared grading head as all Core Four"),
            ("V-threshold", "60 \u2014 the strictest, because rot is dark on golden skin"),
            ("Segmentation", "tuned yellow (15\u201335\u00b0H) \u222a green (25\u201385\u00b0H) mask"),
        ],
    ),
    dict(
        name="Grapes",
        num="04",
        overview=[
            ("What is graded", "bunch-level: shrivel, mold, berry color uniformity, stem health"),
            ("Typical defects", "shriveled berries, powdery mold, uneven color, loose berries"),
            ("Why it is hard", "a bunch is many small objects \u2014 segmentation works on the cluster"),
            ("Grade meaning", "Better = table/export, Good = local market, Reject = juice/discard"),
        ],
        data=[
            ("Format", "dataset-grapes / Better | Good | Reject"),
            ("3 classes", "alphabetical ordering used by TensorFlow"),
            ("Augmentation", "rotation, brightness, zoom, flip, hue"),
            ("Splits", "train / val with held-out validation"),
        ],
        model=[
            ("Weights", "grapes.keras \u2014 27.6 MB MobileNetV2"),
            ("Head", "same shared grading head as all Core Four"),
            ("V-threshold", "45"),
            ("Segmentation", "tuned green (25\u201385\u00b0H) \u222a purple (120\u2013165\u00b0H) mask"),
        ],
    ),
    dict(
        name="Pomegranate",
        num="05",
        overview=[
            ("What is graded", "skin cracking, sunburn, blemishes, dryness, mold"),
            ("Typical defects", "cracks, blackened rind, shrivel, bruising"),
            ("Why it is hard", "deep red color is easy to mask, but cracks look like dark lines to a CNN"),
            ("Grade meaning", "Better = export, Good = retail, Reject = juice/discard"),
        ],
        data=[
            ("Format", "dataset-pomegranate / Better | Good | Reject"),
            ("3 classes", "alphabetical ordering used by TensorFlow"),
            ("Augmentation", "rotation, brightness, zoom, flip, hue"),
            ("Splits", "train / val with held-out validation"),
        ],
        model=[
            ("Weights", "pomegranate.keras \u2014 27.6 MB MobileNetV2"),
            ("Head", "same shared grading head as all Core Four"),
            ("V-threshold", "50"),
            ("Segmentation", "tuned red mask (0\u201310\u00b0H \u222a 160\u2013180\u00b0H)"),
        ],
    ),
]

fruit_slide_no = 12
for f in FRUITS:
    # Overview
    s = content(f"{f['name']} \u2014 Overview", "What HarvestLenz looks for and why it matters.", fruit_slide_no)
    card(s, Inches(0.55), Inches(1.7), Inches(6.0), Inches(4.7), f["name"] + " at a glance", f["overview"], accent=GREEN)
    card(s, Inches(6.85), Inches(1.7), Inches(5.9), Inches(4.7), "Quality classes", [
        ("Better", "high-quality, export-grade appearance"),
        ("Good", "standard retail quality"),
        ("Reject", "damaged / rotten / unmarketable"),
    ], accent=ORANGE)
    txt(s, Inches(0.55), Inches(6.6), Inches(12.2), Inches(0.4),
        "Graded per fruit from a single photo \u2014 no sorting table required.", 15, color=GRAYGREEN, align=PP_ALIGN.CENTER)
    footer(s, fruit_slide_no, dark=False)
    fruit_slide_no += 1

    # Dataset
    s = content(f"{f['name']} \u2014 Dataset & Training", "Three-class quality data, two-phase transfer learning.", fruit_slide_no)
    card(s, Inches(0.55), Inches(1.7), Inches(6.0), Inches(4.7), "Dataset", f["data"], accent=GREEN)
    card(s, Inches(6.85), Inches(1.7), Inches(5.9), Inches(4.7), "Two-phase training", [
        ("Phase 1", "classification head only, lr = 1e-3, ~10 epochs"),
        ("Phase 2", "unfreeze base from layer 100, fine-tune at lr = 1e-4"),
        ("Early stopping", "patience 4\u20135 epochs on validation accuracy"),
        ("Best weights", "checkpoint saved on best val_accuracy"),
    ], accent=ORANGE)
    txt(s, Inches(0.55), Inches(6.6), Inches(12.2), Inches(0.4),
        "The same recipe trains all four Core Four models \u2014 only the data differs.", 15, color=GRAYGREEN, align=PP_ALIGN.CENTER)
    footer(s, fruit_slide_no, dark=False)
    fruit_slide_no += 1

    # Model & behavior
    s = content(f"{f['name']} \u2014 Model & Behavior", "The dedicated weight file and how it behaves.", fruit_slide_no)
    card(s, Inches(0.55), Inches(1.7), Inches(6.0), Inches(4.7), "Dedicated model", f["model"], accent=GREEN)
    card(s, Inches(6.85), Inches(1.7), Inches(5.9), Inches(4.7), "Inference behavior", [
        ("Known fruit", "no classification step \u2014 the dropdown supplies the type"),
        ("Defect fusion", "CNN probability blended with the classical defect score"),
        ("Override", "high defect score forces Reject even if the CNN is confident"),
        ("Confidence gate", "< 0.40 confidence results are discarded"),
    ], accent=ORANGE)
    txt(s, Inches(0.55), Inches(6.6), Inches(12.2), Inches(0.4),
        "Every parameter goes into grading accuracy instead of fruit identification.", 15, color=GRAYGREEN, align=PP_ALIGN.CENTER)
    footer(s, fruit_slide_no, dark=False)
    fruit_slide_no += 1

    # Video demo placeholder
    video_placeholder(f["name"], fruit_slide_no)
    fruit_slide_no += 1

# =============================================================
# SLIDE 28 - End-to-end walkthrough
# =============================================================
s = content("End-to-End Walkthrough", "A mango basket, graded from upload to market call.", 28)
card(s, Inches(0.55), Inches(1.6), Inches(4.1), Inches(4.8), "Input", [
    ("Fruit type", "mango"),
    ("Mode", "basket (is_single = false)"),
    ("Image", "one photo, five mangoes"),
    ("Base price", "Rs. 52.50 / kg"),
], accent=GREEN)
card(s, Inches(4.95), Inches(1.6), Inches(4.1), Inches(4.8), "Per-fruit result", [
    ("Segmented", "5 regions detected, 5 crops extracted"),
    ("Grading", "3 \u00d7 Better (conf. 0.71), 2 \u00d7 Good (conf. 0.83)"),
    ("Per-fruit price", "Better = Rs. 39.38, Good = Rs. 52.50"),
    ("Shelf life", "Better \u2192 ~4 days, Good \u2192 ~7 days"),
], accent=ORANGE)
card(s, Inches(9.35), Inches(1.6), Inches(3.4), Inches(4.8), "Business output", [
    ("Overall grade", "Good (score 85)"),
    ("Market call", "Supermarket \u2014 sell within 48h"),
    ("Basket value", "total \u00d7 1.30 selling estimate"),
    ("Recommendation", "move Better stock first"),
], accent=GREEN)
footer(s, 28, dark=False)

# =============================================================
# SLIDE 29 - API reference
# =============================================================
s = content("API Reference & Session Flow", "Everything is a simple, pollable REST contract.", 29)
rows = [
    ("POST /analyze", "upload + fruit_type + is_single \u2192 session_id"),
    ("GET /analysis/{id}/status", "poll the async job (processing / complete / failed)"),
    ("GET /analysis/{id}", "full result: fruits, summary, pricing, recommendations"),
    ("GET /stats", "dashboard aggregates (scans, fruits, good rate, value)"),
    ("GET /fruits/supported", "the supported crop list"),
    ("GET /health", "model inventory + demo-mode flag"),
]
y = Inches(1.65)
for ep, desc in rows:
    rect(s, Inches(0.55), y, Inches(12.2), Inches(0.78), WHITE)
    txt(s, Inches(0.75), y+Inches(0.06), Inches(4.2), Inches(0.66), ep, 18, bold=True, color=GREEN, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(5.1), y+Inches(0.06), Inches(7.5), Inches(0.66), desc, 15, color=GRAYGREEN, anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.88)
txt(s, Inches(0.55), Inches(6.75), Inches(12.2), Inches(0.4),
    "Flow:  Upload \u2192 session_id \u2192 poll status \u2192 fetch results \u2192 render report", 16, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
footer(s, 29, dark=False)

# =============================================================
# SLIDE 30 - Tech stack
# =============================================================
s = content("Tech Stack", "The tools behind the grade.", 30)
st = [
    ("FastAPI", "async backend, session-based API, background jobs"),
    ("OpenCV", "classical segmentation \u2014 no ML detector"),
    ("TensorFlow / Keras", "MobileNetV2 transfer learning, two-phase fine-tune"),
    ("SQLAlchemy (async)", "session + per-fruit persistence"),
    ("SQLite", "zero-infrastructure local database"),
    ("Vanilla HTML/JS", "dashboard and scan pages, no build tooling"),
]
x, y = Inches(0.55), Inches(1.75)
for i, (t, d) in enumerate(st):
    cx = x + (i % 2) * Inches(6.2)
    cy = y + (i // 2) * Inches(1.65)
    card(s, cx, cy, Inches(6.0), Inches(1.45), t, d, accent=GREEN if i % 2 == 0 else ORANGE)
txt(s, Inches(0.55), Inches(6.7), Inches(12.2), Inches(0.4),
    "Philosophy: deterministic rules and classical vision where they win, deep learning only where it earns its keep.", 16, bold=True, color=GRAYGREEN, align=PP_ALIGN.CENTER)
footer(s, 30, dark=False)

# =============================================================
# SLIDE 31 - Roadmap & closing
# =============================================================
s = new_slide(DARK)
rect(s, 0, 0, SW, Inches(0.18), GREEN)
txt(s, Inches(0.6), Inches(0.7), Inches(12.0), Inches(0.7), "Roadmap & Closing", 40, bold=True, color=WHITE, font=FONT_H)
rect(s, Inches(0.65), Inches(1.45), Inches(0.9), Inches(0.06), ORANGE)
bullets(s, Inches(0.6), Inches(1.8), Inches(6.6), Inches(4.4), [
    ("Validate per-crop accuracy", "formal eval on held-out splits, per class"),
    ("Harden segmentation", "augment training with cluttered-basket captures"),
    ("More tuned crops", "extend tuned masks beyond the Core Four"),
    ("Edge deployment", "quantized TFLite (INT8/FP16) for on-device inference"),
    ("Mobile app", "point, shoot, and get a grade report on the phone"),
], size=18, color=RGBColor(0xDD, 0xE0, 0xD5))
card(s, Inches(7.5), Inches(1.8), Inches(5.2), Inches(4.4), "The Core Four today", [
    ("Mango", "MobileNetV2 \u00b7 V-threshold 40 \u00b7 generic mask"),
    ("Pineapple", "MobileNetV2 \u00b7 V-threshold 60 \u00b7 tuned mask"),
    ("Grapes", "MobileNetV2 \u00b7 V-threshold 45 \u00b7 tuned mask"),
    ("Pomegranate", "MobileNetV2 \u00b7 V-threshold 50 \u00b7 tuned mask"),
], bg=WHITE)
txt(s, Inches(0.6), Inches(6.4), Inches(12.2), Inches(0.6), "One photo in.  Grade, shelf life, and a market call out. \u2014 Thank you.", 24, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
footer(s, 31, dark=True)

# ---------- save ----------
out = os.environ.get("HL_PPTX_OUT", r"C:\Users\madha\OneDrive\Desktop\HarvestLenz_CoreFour_Presentation.pptx")
prs.save(out)
print("Saved:", out)
print("Slides:", len(prs.slides.__iter__.__self__._sldIdLst))
